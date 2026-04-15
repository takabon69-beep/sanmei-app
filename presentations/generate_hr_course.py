from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime
import os

# 16:9 のサイズ定義 (13.333 x 7.5 inches)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

def apply_branding(slide, prs, page_num):
    """ブランディング（ロゴ、フッター、ページ番号）を適用する (16:9対応)"""
    # --- ロゴの配置 (左上) ---
    logo_path = os.path.join("presentations", "brand", "logo.png")
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, Inches(0.4), Inches(0.3), height=Inches(0.8))
    
    # --- コピーライト (フッター中央) ---
    copyright_text = "(C) 2026 株式会社村上経営研究所"
    footer_box = slide.shapes.add_textbox(Inches(0), Inches(6.8), SLIDE_WIDTH, Inches(0.4))
    ftf = footer_box.text_frame
    p = ftf.paragraphs[0]
    p.text = copyright_text
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(128, 128, 128)
    p.alignment = PP_ALIGN.CENTER
    
    # --- ページ番号 (右下) ---
    page_num_box = slide.shapes.add_textbox(SLIDE_WIDTH - Inches(1.2), Inches(6.8), Inches(0.8), Inches(0.4))
    ptf = page_num_box.text_frame
    p = ptf.paragraphs[0]
    p.text = str(page_num)
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(128, 128, 128)
    p.alignment = PP_ALIGN.RIGHT

def add_mascot_with_bubble(slide, bubble_text=""):
    """マスコット（Nanobanana）と吹き出しを配置する"""
    mascot_path = os.path.join("presentations", "brand", "mascot", "mascot_main.png")
    if os.path.exists(mascot_path):
        # マスコットを右下に配置
        m_width = Inches(2.8)
        m_height = Inches(2.8)
        m_left = SLIDE_WIDTH - m_width - Inches(0.5)
        m_top = SLIDE_HEIGHT - m_height - Inches(0.8)
        slide.shapes.add_picture(mascot_path, m_left, m_top, width=m_width)
        
        # 吹き出しの追加
        if bubble_text:
            # 吹き出しの種類を WEDGE_ROUND_RECT_CALLOUT (定数値 109) で指定
            # python-pptx のバージョンにより定数名が異なる場合があるため整数値を使用
            bubble = slide.shapes.add_shape(
                109, # MSO_SHAPE.WEDGE_ROUND_RECT_CALLOUT
                m_left - Inches(3.5), m_top - Inches(0.5), Inches(4.0), Inches(1.5)
            )
            bubble.fill.solid()
            bubble.fill.fore_color.rgb = RGBColor(255, 255, 230) # 薄い黄色
            bubble.line.color.rgb = RGBColor(150, 150, 0)
            
            tf = bubble.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = bubble_text
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 0, 0)
            p.alignment = PP_ALIGN.CENTER

def add_title_slide(prs, title, subtitle, info):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    apply_branding(slide, prs, 1)
    
    # タイトルスライドには大きめのマスコット
    mascot_path = os.path.join("presentations", "brand", "mascot", "mascot_main.png")
    if os.path.exists(mascot_path):
        slide.shapes.add_picture(mascot_path, SLIDE_WIDTH - Inches(4.5), Inches(2.2), width=Inches(4.0))
    
    # メインタイトル (左寄せ)
    text_box = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(8), Inches(2.0))
    tf = text_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    
    # サブタイトル
    p = tf.add_paragraph()
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(40, 40, 40)
    p.alignment = PP_ALIGN.LEFT
    
    # 属性
    info_box = slide.shapes.add_textbox(Inches(1), Inches(6.0), Inches(8), Inches(0.5))
    itf = info_box.text_frame
    p = itf.paragraphs[0]
    p.text = info
    p.font.size = Pt(14)
    p.alignment = PP_ALIGN.LEFT
    
    return slide

def add_content_slide(prs, title, items, page_num, notes="", bubble_text=""):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    apply_branding(slide, prs, page_num)
    add_mascot_with_bubble(slide, bubble_text)
    
    # 見出し
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), SLIDE_WIDTH - Inches(1), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 112, 192)
    
    # コンテンツ
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.5), Inches(4.5))
    ctf = content_box.text_frame
    ctf.word_wrap = True
    for item in items:
        p = ctf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(22)
        p.space_after = Pt(12)
    
    # ノート
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    
    return slide

def add_section_slide(prs, title1, title2, page_num, notes=""):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    apply_branding(slide, prs, page_num)
    
    # 中央に大きく表示
    box = slide.shapes.add_textbox(Inches(1), Inches(2.5), SLIDE_WIDTH - Inches(2), Inches(2.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title1
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(84, 130, 53)
    p.alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = title2
    p.font.size = Pt(54)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # セクションスライドにはマスコットを左側に配置
    mascot_path = os.path.join("presentations", "brand", "mascot", "mascot_main.png")
    if os.path.exists(mascot_path):
        slide.shapes.add_picture(mascot_path, Inches(0.5), SLIDE_HEIGHT - Inches(3.5), width=Inches(3.0))

    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    
    return slide

def generate():
    prs = Presentation()
    # 16:9 へのサイズ変更
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    
    # Slide 01: Title
    add_title_slide(prs, "店長のための人材育成講座", "― まず自分を育てる、そして部下を育てる ―", 
                    "対象：小売業 店長 ／ 形式：ワークショップ（2時間）")
    
    # Slide 02
    add_content_slide(prs, "部下を育てる前に", 
                      ["あなた自身は、育ち続けていますか？", 
                       "・この講座は「部下育成のテクニック」を教えません", 
                       "・店長自身の「在り方」を問い直す2時間です",
                       "自分を育てる → 感情を整える → 伝わる言葉 → 部下が育つ"], 2,
                      notes="テクニックではなく在り方を強調。最後は矢印のシンプルな図としてイメージ。",
                      bubble_text="まずは店長自身の「在り方」が大切ですよ！")
    
    # Slide 03
    add_content_slide(prs, "隣の人と話してみましょう（3分）", 
                      ["「今、店長として一番しんどいことは何ですか？」", 
                       "・正解はありません", 
                       "・まず自分が話す、次に聞く"], 3,
                      notes="心理的安全を整える一言を添える。一言でいいですよと気軽に促す。",
                      bubble_text="お隣の方と気軽に話してみましょう！")
    
    # Slide 04
    add_section_slide(prs, "第一部① 学ぶ姿勢・成長マインドセット", "「成長し続ける店長」と「止まってしまう店長」", 4,
                      notes="才能でも環境でもないことを断言。キーワードはマインドセット。")
    
    # Slide 05
    add_content_slide(prs, "固定型 vs 成長型　どちらにありますか？", 
                      ["【固定型】「向いてない」「能力は変わらない」「忙しいから無理」", 
                       "【成長型】「失敗から学ぶ」「盗めるものはないか」「5分でもやる」",
                       "マインドセットは「気づいた瞬間」から変えられる"], 5,
                      notes="自己診断として投げかけ、切り替えられることが重要であると〆る。",
                      bubble_text="「気づいた瞬間」から成長は始まります！")
    
    # Slide 06
    add_content_slide(prs, "方眼ノートに書いてみましょう（5分）", 
                      ["① 自分が最後に「本気で学んだ」のはいつですか？", 
                       "② そのきっかけは？その後どう変わりましたか？",
                       "③ 「忙しくて学べない」とき、何が優先されていますか？"], 6,
                      notes="静かな時間をつくる。書けない人にはなぜ書けないかを書くよう促す。",
                      bubble_text="自分自身を振り返る時間です。")
    
    # Slide 07
    add_content_slide(prs, "店長として自分をどれくらい知っていますか？", 
                      ["① 得意なこと（強み）→ どんな場面で発揮？", 
                       "② 苦手なこと（弱み）→ それを補うために？",
                       "③ 大切なこと（価値観）→ 譲れないものは？"], 7,
                      notes="自己認識の精度が関わり方を決める。ジョハリの窓に触れても良い。",
                      bubble_text="自分の「盲点」を知ることが部下理解のヒントです！")
    
    # Slide 08
    add_content_slide(prs, "強みを相手の言葉で伝え合う（各3分）", 
                      ["「私が感じるあなたの強みは、〇〇なところです」", 
                       "・自己申告ではなく、他者からの言葉を受け取る",
                       "・自分の見えていない面を知ることが部下理解の入口"], 8,
                      notes="ズレがあることが学びになる。第二部への橋渡し。",
                      bubble_text="他人の目線は新しい発見がありますね！")
    
    # Slide 09
    add_section_slide(prs, "第二部① 感情のクセを知り、整える", "店長の感情が売場をつくる", 9,
                      notes="感情伝染（エモーショナルコンテイジョン）について説明。")
    
    # Slide 10
    add_content_slide(prs, "「〜すべき」の裏側にあるもの（7分）", 
                      ["① 最近、イライラした場面を3つ挙げてください", 
                       "② そのとき期待していた「〜すべき」は何ですか？",
                       "③ その「べき」は誰から学んだものですか？",
                       "怒り ＝ 「べき思考」が裏切られたサイン"], 10,
                      notes="誰にも見せなくていい。怒りの正体（べき思考）を言語化。",
                      bubble_text="イライラは「べき」の裏返しなんです。")
    
    # Slide 11
    add_content_slide(prs, "感情コントロールは「変換」の技術", 
                      ["・6秒ルール：ピークが過ぎるまで反応を待つ", 
                       "・「なんで！」を「してほしかった（要望）」へ",
                       "× また遅刻か！ → ○ 時間通りに来てほしかった。何かあった？",
                       "× 何度言えば！ → ○ 手順をもう一度一緒に確認したい"], 11,
                      notes="我慢ではなく変換。次のコミュニケーションパートへ橋渡し。",
                      bubble_text="怒りを「要望」に変えてみましょう！")
    
    # Slide 12
    add_section_slide(prs, "第二部② 伝わるコミュニケーションの設計", "「指示したつもり」の落とし穴", 12,
                      notes="指示 vs 対話。相手が自分の言葉で動き出す状態を目指す。")
    
    # Slide 13
    add_content_slide(prs, "遅刻が続く部下への関わり方を変える", 
                      ["① 命令型：「またか。何度言えばわかる」", 
                       "② 承認型：「最近頑張ってるね。遅刻だけ改善しよう」",
                       "③ 問いかけ型：「最近続いてるけど、何か困ってる？」",
                       "相手に一番届くのはどのタイプ？"], 13,
                      notes="実際に声に出して体感させる。どのタイプが届くかを気づかせる。",
                      bubble_text="相手の背景を聞くスタンスが大切です。")
    
    # Slide 14
    add_content_slide(prs, "承認の4段階", 
                      ["1. 【存在承認】「おはよう」「ありがとう」", 
                       "2. 【行動承認】「〇〇してくれたね」",
                       "3. 【成果承認】「目標達成！よくやった」",
                       "4. 【成長承認】「先月より確実に成長してるね」",
                       "第4段階まで使い、人を育てる"], 14,
                      notes="存在・行動・成果・成長。行動やプロセスに注目するのがポイント。",
                      bubble_text="「成長承認」ができる店長を目指しましょう！")
    
    # Slide 15
    add_section_slide(prs, "第三部 育つ場の設計と明日からの行動", "「育てる」のではなく「育つ場をつくる」", 15,
                      notes="店長の設計力が試される最終パート。")
    
    # Slide 16
    add_content_slide(prs, "明日、まず誰の、何を認めますか？", 
                      ["行動宣言：明日、【　】に対し、【　】をする", 
                       "・具体的、かつ小さな一歩に絞る",
                       "・「部下が育つ店は、店長が一番学んでいる店」",
                       "ペアになり「それ、いいですね」と承認し合う"], 16,
                      notes="今日学んだ承認をその場で実践する。問いを最後に残して終わる。",
                      bubble_text="まずは小さな一歩から始めてみましょう！")

    # 保存 (v3 として保存)
    date_str = datetime.now().strftime("%Y%m%d")
    title_filename = "店長の人材育成講座"
    filename = f"{date_str}_{title_filename}_v3.pptx"
    
    output_dir = os.path.join("presentations", "outputs")
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
        
    output_path = os.path.join(output_dir, filename)
    prs.save(output_path)
    print(f"スライドを生成しました: {output_path}")

if __name__ == "__main__":
    generate()
