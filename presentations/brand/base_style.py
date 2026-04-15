from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_base_slide(title_text, content_list, logo_path=None, copyright_text="(C) 2026 株式会社村上経営研究所"):
    # プレゼンテーションの作成
    prs = Presentation()
    
    # スライドの追加 (レイアウトは「白紙」を使用)
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # --- ロゴの配置 (左上) ---
    if logo_path and os.path.exists(logo_path):
        # 縦横比を維持しながら配置 (例: 左から0.5インチ、上から0.3インチ、高さ0.8インチ)
        slide.shapes.add_picture(logo_path, Inches(0.5), Inches(0.3), height=Inches(0.8))
    
    # --- タイトルの配置 ---
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)
    
    # --- 内容の配置 (箇条書き) ---
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(8.5), Inches(4.5))
    ctf = content_box.text_frame
    ctf.word_wrap = True
    for item in content_list:
        p = ctf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(24)
        p.space_after = Pt(10)
    
    # --- コピーライト (フッター) ---
    # 横幅をスライド全体（10インチ）に広げて中央揃えにする
    footer_box = slide.shapes.add_textbox(Inches(0), Inches(6.8), Inches(10), Inches(0.5))
    ftf = footer_box.text_frame
    ftf.word_wrap = True
    p = ftf.paragraphs[0]
    p.text = copyright_text
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(128, 128, 128) # グレー
    p.alignment = PP_ALIGN.CENTER
    
    # --- ページ番号 (右下) ---
    page_num_box = slide.shapes.add_textbox(Inches(8.5), Inches(6.8), Inches(1), Inches(0.5))
    ptf = page_num_box.text_frame
    p = ptf.paragraphs[0]
    p.text = "1" # サンプルのため 1 をセット
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(128, 128, 128)
    p.alignment = PP_ALIGN.RIGHT
    
    return prs

if __name__ == "__main__":
    from datetime import datetime
    
    # プレゼン情報
    title = "プレゼンテーション設計の標準化"
    date_str = datetime.now().strftime("%Y%m%d")
    filename = f"{date_str}_{title}.pptx"
    
    # テスト実行
    logo_file = os.path.join("presentations", "brand", "logo.png")
    prs_test = create_base_slide(
        title,
        ["ロゴは左上に配置", "フッターにコピーライト（中央・2026年）", "右下にページ番号を表示", "構成案から自動生成が可能"]
    )
    
    # 保存先 (outputs/ フォルダへ)
    output_dir = os.path.join("presentations", "outputs")
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
        
    output_path = os.path.join(output_dir, filename)
    prs_test.save(output_path)
    print(f"サンプルスライドを保存しました: {output_path}")
